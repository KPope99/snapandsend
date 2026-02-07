#!/usr/bin/env python3
"""
Generate SnapAndSend & Incident Response Proposal PowerPoint
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# Create presentation
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Color scheme
PRIMARY_GREEN = RGBColor(16, 185, 129)  # Emerald-500
DARK_GREEN = RGBColor(6, 95, 70)        # Emerald-800
LIGHT_GREEN = RGBColor(209, 250, 229)   # Emerald-100
DARK_GRAY = RGBColor(31, 41, 55)        # Gray-800
LIGHT_GRAY = RGBColor(107, 114, 128)    # Gray-500
WHITE = RGBColor(255, 255, 255)
BLUE = RGBColor(59, 130, 246)           # Blue-500

def add_title_slide(title, subtitle):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = DARK_GREEN
    shape.line.fill.background()

    # Accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(3.2), prs.slide_width, Inches(0.1))
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY_GREEN
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(12.333), Inches(1.2))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(12.333), Inches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = subtitle
    p.font.size = Pt(24)
    p.font.color.rgb = LIGHT_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12.333), Inches(0.5))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = "© Tech84 | Community Incident Reporting Platform"
    p.font.size = Pt(14)
    p.font.color.rgb = LIGHT_GREEN
    p.alignment = PP_ALIGN.CENTER

    return slide

def add_section_slide(title):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Left accent bar
    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.3), prs.slide_height)
    accent.fill.solid()
    accent.fill.fore_color.rgb = PRIMARY_GREEN
    accent.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(3), Inches(11.533), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.LEFT

    return slide

def add_content_slide(title, content_items, icon_text=None):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Content
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12.333), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = f"• {item['title']}"
            p.font.size = Pt(22)
            p.font.bold = True
            p.font.color.rgb = DARK_GREEN
            p.space_after = Pt(4)

            if 'desc' in item:
                p2 = tf.add_paragraph()
                p2.text = f"   {item['desc']}"
                p2.font.size = Pt(18)
                p2.font.color.rgb = LIGHT_GRAY
                p2.space_after = Pt(16)
        else:
            p.text = f"• {item}"
            p.font.size = Pt(20)
            p.font.color.rgb = DARK_GRAY
            p.space_after = Pt(12)

    # Footer
    footer = slide.shapes.add_textbox(Inches(11.5), Inches(7), Inches(1.5), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "© Tech84"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT

    return slide

def add_two_column_slide(title, left_title, left_items, right_title, right_items):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Left column box
    left_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.4), Inches(1.5), Inches(6), Inches(5.5))
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = LIGHT_GREEN
    left_box.line.color.rgb = PRIMARY_GREEN

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.7), Inches(5.6), Inches(0.6))
    tf = left_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = left_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = DARK_GREEN
    p.alignment = PP_ALIGN.CENTER

    # Left column content
    left_content = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(5.4), Inches(4.4))
    tf = left_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(left_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # Right column box
    right_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.9), Inches(1.5), Inches(6), Inches(5.5))
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = RGBColor(219, 234, 254)  # Blue-100
    right_box.line.color.rgb = BLUE

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(7.1), Inches(1.7), Inches(5.6), Inches(0.6))
    tf = right_title_box.text_frame
    p = tf.paragraphs[0]
    p.text = right_title
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = RGBColor(30, 64, 175)  # Blue-800
    p.alignment = PP_ALIGN.CENTER

    # Right column content
    right_content = slide.shapes.add_textbox(Inches(7.2), Inches(2.4), Inches(5.4), Inches(4.4))
    tf = right_content.text_frame
    tf.word_wrap = True
    for i, item in enumerate(right_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(8)

    # Footer
    footer = slide.shapes.add_textbox(Inches(11.5), Inches(7), Inches(1.5), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "© Tech84"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT

    return slide

def add_workflow_slide(title, steps):
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = DARK_GREEN
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12.333), Inches(0.7))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = WHITE

    # Steps
    step_width = (prs.slide_width - Inches(1)) / len(steps)
    colors = [PRIMARY_GREEN, BLUE, RGBColor(245, 158, 11), RGBColor(239, 68, 68), RGBColor(139, 92, 246)]

    for i, step in enumerate(steps):
        x = Inches(0.5) + (step_width * i)

        # Step box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x + Inches(0.1), Inches(1.8), step_width - Inches(0.2), Inches(4.5))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(249, 250, 251)  # Gray-50
        box.line.color.rgb = colors[i % len(colors)]

        # Step number
        num_box = slide.shapes.add_shape(MSO_SHAPE.OVAL, x + step_width/2 - Inches(0.3), Inches(2), Inches(0.6), Inches(0.6))
        num_box.fill.solid()
        num_box.fill.fore_color.rgb = colors[i % len(colors)]
        num_box.line.fill.background()

        num_text = slide.shapes.add_textbox(x + step_width/2 - Inches(0.3), Inches(2.1), Inches(0.6), Inches(0.5))
        tf = num_text.text_frame
        p = tf.paragraphs[0]
        p.text = str(i + 1)
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.alignment = PP_ALIGN.CENTER

        # Step title
        step_title = slide.shapes.add_textbox(x + Inches(0.2), Inches(2.8), step_width - Inches(0.4), Inches(0.6))
        tf = step_title.text_frame
        p = tf.paragraphs[0]
        p.text = step['title']
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = DARK_GRAY
        p.alignment = PP_ALIGN.CENTER

        # Step description
        step_desc = slide.shapes.add_textbox(x + Inches(0.2), Inches(3.4), step_width - Inches(0.4), Inches(2.5))
        tf = step_desc.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = step['desc']
        p.font.size = Pt(12)
        p.font.color.rgb = LIGHT_GRAY
        p.alignment = PP_ALIGN.CENTER

        # Arrow between steps
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x + step_width - Inches(0.15), Inches(4), Inches(0.3), Inches(0.3))
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = LIGHT_GRAY
            arrow.line.fill.background()

    # Footer
    footer = slide.shapes.add_textbox(Inches(11.5), Inches(7), Inches(1.5), Inches(0.4))
    tf = footer.text_frame
    p = tf.paragraphs[0]
    p.text = "© Tech84"
    p.font.size = Pt(10)
    p.font.color.rgb = LIGHT_GRAY
    p.alignment = PP_ALIGN.RIGHT

    return slide

# ============= CREATE SLIDES =============

# Slide 1: Title
add_title_slide(
    "SnapAndSend & Incident Response",
    "A Community-Powered Incident Reporting & Resolution Platform"
)

# Slide 2: Executive Summary
add_section_slide("Executive Summary")

# Slide 3: The Problem
add_content_slide(
    "The Problem We're Solving",
    [
        {"title": "Delayed Incident Reporting", "desc": "Traditional reporting methods are slow, bureaucratic, and often ignored"},
        {"title": "Lack of Transparency", "desc": "Citizens don't know if their reports are being addressed or the status of resolution"},
        {"title": "Unverified Reports", "desc": "Authorities struggle to prioritize genuine incidents from false reports"},
        {"title": "No Accountability", "desc": "No tracking of response times or resolution effectiveness"},
        {"title": "Communication Gap", "desc": "Disconnect between community members and responding authorities"}
    ]
)

# Slide 4: Our Solution
add_content_slide(
    "Our Solution: Two Integrated Applications",
    [
        {"title": "SnapAndSend (Community App)", "desc": "Mobile-first PWA for citizens to report incidents with photos, location, and real-time verification"},
        {"title": "Incident Response (Authority Dashboard)", "desc": "Comprehensive dashboard for police/authorities to manage, investigate, and resolve reported incidents"},
        {"title": "Seamless Integration", "desc": "Real-time sync via webhooks and External API for third-party systems integration"},
        {"title": "AI-Powered Analysis", "desc": "Automatic incident categorization and duplicate detection using computer vision"}
    ]
)

# Slide 5: Two Apps Overview
add_two_column_slide(
    "Platform Overview",
    "SnapAndSend (Citizens)",
    [
        "Report incidents with photo evidence",
        "GPS-based location tagging",
        "AI-powered category detection",
        "Community verification system",
        "Track report status in real-time",
        "View nearby incidents on map",
        "Receive resolution notifications",
        "Works offline as PWA"
    ],
    "Incident Response (Authorities)",
    [
        "Centralized incident dashboard",
        "Real-time incident feed (SSE)",
        "Status management workflow",
        "Resolution with evidence upload",
        "Timeline tracking & audit logs",
        "Statistics and analytics",
        "External API for integrations",
        "Webhook notifications"
    ]
)

# Slide 6: How It Works
add_workflow_slide(
    "How The Platform Works",
    [
        {"title": "Report", "desc": "Citizen captures photo of incident. AI analyzes and suggests category."},
        {"title": "Verify", "desc": "Nearby users verify the report. Duplicates auto-merge as verifications."},
        {"title": "Investigate", "desc": "Authorities receive alert, review incident, begin investigation."},
        {"title": "Resolve", "desc": "Authority uploads evidence, adds notes, marks resolved."},
        {"title": "Notify", "desc": "Reporter and verifiers notified. Timeline logged for transparency."}
    ]
)

# Slide 7: SnapAndSend Features
add_content_slide(
    "SnapAndSend - Core Features",
    [
        {"title": "Smart Photo Capture", "desc": "Camera integration with AI-powered incident analysis and auto-categorization"},
        {"title": "GPS Location Tracking", "desc": "Automatic location detection with manual override option for accurate placement"},
        {"title": "Community Verification", "desc": "Nearby users can verify incidents (within 500m), boosting credibility"},
        {"title": "Duplicate Detection", "desc": "Automatic merging of similar incidents within 200m as verifications"},
        {"title": "Real-time Map View", "desc": "Interactive map showing all nearby incidents with status indicators"},
        {"title": "Session Management", "desc": "30-minute timeout for security with activity tracking"},
        {"title": "Progressive Web App", "desc": "Install on any device, works offline, push notifications"}
    ]
)

# Slide 8: Incident Response Features
add_content_slide(
    "Incident Response - Core Features",
    [
        {"title": "Live Dashboard", "desc": "Real-time incident feed with Server-Sent Events (SSE) for instant updates"},
        {"title": "Status Workflow", "desc": "Pending → Investigating → Resolved with timestamp logging"},
        {"title": "Evidence Management", "desc": "Mandatory evidence upload and remediation notes before resolution"},
        {"title": "Timeline Tracking", "desc": "Full audit trail: report time, investigation start, resolution time"},
        {"title": "External API", "desc": "RESTful API with API key auth for third-party system integration"},
        {"title": "Webhook Support", "desc": "Real-time notifications to external systems on incident events"},
        {"title": "Statistics Dashboard", "desc": "Analytics on incident types, response times, resolution rates"}
    ]
)

# Slide 9: Incident Categories
add_content_slide(
    "Supported Incident Categories",
    [
        {"title": "Infrastructure", "desc": "Potholes, road damage, streetlight outages, drainage issues, damaged signage"},
        {"title": "Environmental", "desc": "Illegal dumping, garbage overflow, flooding, pollution"},
        {"title": "Public Safety", "desc": "Vandalism, robbery, assault, suspicious activity"},
        {"title": "Traffic", "desc": "Traffic light malfunction, road blockages, accidents"},
        {"title": "AI-Detected Categories", "desc": "System automatically detects and suggests new categories from image analysis"},
        {"title": "Custom Categories", "desc": "Authorities can define region-specific incident types"}
    ]
)

# Slide 10: Technical Architecture
add_content_slide(
    "Technical Architecture",
    [
        {"title": "Frontend", "desc": "React + TypeScript + Vite, TailwindCSS, Leaflet Maps, PWA-ready"},
        {"title": "Backend", "desc": "Node.js + Express, Prisma ORM, SQLite (dev) / PostgreSQL (prod)"},
        {"title": "AI Integration", "desc": "OpenAI Vision API (GPT-4o) for image analysis and categorization"},
        {"title": "Real-time", "desc": "Server-Sent Events (SSE) for live updates, webhooks for integrations"},
        {"title": "Security", "desc": "JWT authentication, API key validation, session timeout, HTTPS"},
        {"title": "Storage", "desc": "Local file storage with S3-compatible cloud storage option"}
    ]
)

# Slide 11: Community Benefits
add_content_slide(
    "Benefits to the Community",
    [
        {"title": "Empowered Citizens", "desc": "Easy way to report issues and track resolution - voice is heard"},
        {"title": "Faster Response", "desc": "Real-time alerts mean quicker authority response to critical incidents"},
        {"title": "Transparency", "desc": "Full visibility into incident status, timeline, and resolution evidence"},
        {"title": "Accountability", "desc": "Audit trails ensure authorities are held responsible for timely resolution"},
        {"title": "Community Trust", "desc": "Verified reports from multiple citizens increase credibility"},
        {"title": "Safer Neighborhoods", "desc": "Proactive incident reporting prevents escalation and improves safety"},
        {"title": "Data-Driven Decisions", "desc": "Analytics help identify problem areas for targeted improvements"}
    ]
)

# Slide 12: Benefits to Authorities
add_content_slide(
    "Benefits to Authorities",
    [
        {"title": "Centralized Management", "desc": "Single dashboard for all community-reported incidents"},
        {"title": "Prioritized Response", "desc": "Verified incidents with multiple confirmations get priority"},
        {"title": "Reduced False Reports", "desc": "Community verification and AI analysis filter out invalid reports"},
        {"title": "Evidence Collection", "desc": "Photo evidence from multiple angles and locations"},
        {"title": "Performance Metrics", "desc": "Track response times, resolution rates, and team performance"},
        {"title": "Integration Ready", "desc": "API and webhooks connect with existing dispatch/CAD systems"},
        {"title": "Public Relations", "desc": "Demonstrate responsiveness and transparency to citizens"}
    ]
)

# Slide 13: API Integration
add_content_slide(
    "External API & Integration",
    [
        {"title": "RESTful API Endpoints", "desc": "GET /incidents, GET /incidents/:id, PATCH /incidents/:id/status, GET /stats"},
        {"title": "API Key Authentication", "desc": "Secure access with X-API-Key header, partner management"},
        {"title": "Webhook Events", "desc": "incident.created, incident.verified, incident.status_changed, incident.resolved"},
        {"title": "Status Management", "desc": "External systems can update status: pending → investigating → resolved"},
        {"title": "Location Filtering", "desc": "Query by lat/lng/radius for jurisdiction-based filtering"},
        {"title": "Pagination Support", "desc": "Limit/offset parameters for handling large datasets"}
    ]
)

# Slide 14: Security Features
add_content_slide(
    "Security & Privacy",
    [
        {"title": "User Authentication", "desc": "Secure registration/login with password hashing (bcrypt)"},
        {"title": "Session Management", "desc": "30-minute inactivity timeout with secure token handling"},
        {"title": "API Security", "desc": "API key validation, rate limiting, partner access logging"},
        {"title": "Data Privacy", "desc": "Optional anonymous reporting, minimal PII collection"},
        {"title": "Audit Logging", "desc": "All status changes logged with timestamp and actor"},
        {"title": "HTTPS Encryption", "desc": "All data transmitted over secure encrypted connections"}
    ]
)

# Slide 15: Implementation Roadmap
add_workflow_slide(
    "Implementation Roadmap",
    [
        {"title": "Phase 1\nFoundation", "desc": "Core reporting, map view, basic auth, incident management"},
        {"title": "Phase 2\nAI & Verification", "desc": "AI categorization, community verification, duplicate detection"},
        {"title": "Phase 3\nAuthority Tools", "desc": "Dashboard, status workflow, evidence upload, timeline"},
        {"title": "Phase 4\nIntegration", "desc": "External API, webhooks, third-party system connections"},
        {"title": "Phase 5\nScale", "desc": "Analytics, multi-region, mobile apps, advanced reporting"}
    ]
)

# Slide 16: Use Cases
add_content_slide(
    "Real-World Use Cases",
    [
        {"title": "Municipal Services", "desc": "City councils receive and track infrastructure repair requests"},
        {"title": "Police Departments", "desc": "Crime reporting with verified community witnesses"},
        {"title": "Emergency Services", "desc": "Flood, fire, or accident reporting with real-time location"},
        {"title": "Environmental Agencies", "desc": "Track illegal dumping and pollution incidents"},
        {"title": "Neighborhood Watch", "desc": "Community-organized safety monitoring and reporting"},
        {"title": "Utility Companies", "desc": "Report outages, damaged infrastructure, safety hazards"}
    ]
)

# Slide 17: Contact/Next Steps
add_title_slide(
    "Ready to Transform\nCommunity Safety?",
    "Contact: Tech84 | Let's discuss implementation for your region"
)

# Save presentation
output_path = "/Users/olahpope/snapandsend/SnapAndSend_Proposal.pptx"
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
